import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def add_slide(prs, title, content_bullets):
    slide_layout = prs.slide_layouts[1] # Title and content
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for idx, (bullet_text, font_size, is_bold) in enumerate(content_bullets):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = bullet_text
        p.font.size = Pt(font_size)
        p.font.bold = is_bold
        p.font.color.rgb = RGBColor(0, 0, 0)

def main():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Neurite Detection Pipeline Algorithm"
    subtitle.text = "End-to-End Processing & Core Routines"

    # Slide 1
    add_slide(prs, "1. Soma Detection (488 Channel)", [
        ("Goal: Isolate cell bodies to prevent internal skeletonization", 24, True),
        ("Core Routines & Technical Names:", 20, True),
        (" - Parameter Tuning: Gamma & Contrast (np.clip, np.power)", 18, False),
        (" - Smoothing: cupyx.scipy.ndimage.gaussian_filter", 18, False),
        (" - Thresholding: cupy.percentile (> 98.5%)", 18, False),
        (" - Morphological Clean-up: cp_ndi.binary_erosion / binary_dilation", 18, False),
        (" - Connected Components: cupyx.scipy.ndimage.label", 18, False),
        (" - Spatial Filtering: > 3000 voxels (cp.bincount)", 18, False),
        (" - Centroid Extraction: scipy.ndimage.center_of_mass", 18, False)
    ])

    # Slide 2
    add_slide(prs, "2. Neurite Manifold & Skeletonization", [
        ("Goal: Extract 3D neurite manifolds and skeletonize to 1-pixel widths", 24, True),
        ("Core Routines & Technical Names:", 20, True),
        (" - Vesselness Filtering: process_tile_frangi_gpu (Custom Analytical GPU Frangi)", 18, False),
        (" - Eigenvalue Calculation: eigh_3x3_analytical (Cardano's method on GPU)", 18, False),
        (" - Manifold Hysteresis: cp_ndi.label & High/Low dual thresholds", 18, False),
        (" - Soma Subtraction: Boolean Masking (gpu_binary & ~gpu_smask)", 18, False),
        (" - Fragment Removal: Topological > 3000 voxel constraints", 18, False),
        (" - 3D Skeletonization: skimage.morphology.skeletonize (Multi-threaded CPU chunks)", 18, False)
    ])

    # Slide 3
    add_slide(prs, "3. Barcode Detection & Filtering (555/640)", [
        ("Goal: Detect barcodes and rigorously prune false positives", 24, True),
        ("Core Routines & Technical Names:", 20, True),
        (" - Barcode Spot Detection: detect_barcodes (cp.percentile > 99.9%)", 18, False),
        (" - Component Extraction: cupyx.scipy.ndimage.label & center_of_mass", 18, False),
        (" - Spatial Buffer Creation: cp_ndi.binary_dilation (+-5 voxels around skeletons)", 18, False),
        (" - Structure Filtering: filter_structures (Boolean Intersections)", 18, False),
        (" - Connectivity Pruning: filter_mask_by_barcodes (Isolating validated structures)", 18, False)
    ])

    # Slide 4
    add_slide(prs, "4. Graph Network Tracing & Export", [
        ("Goal: Convert raster skeletons into mathematical network graphs", 24, True),
        ("Core Routines & Technical Names:", 20, True),
        (" - Subgraph Identification: scipy.ndimage.label (299,000+ independent components)", 18, False),
        (" - Path Tracing: export_graphs (Iterative topological traversal)", 18, False),
        (" - Graph Structuring: Constructing JSON CW Complex", 18, False),
        (" - Morphological Standard: SWC file generation with physical scaling", 18, False),
        (" - Proofreading Integration: napari.Viewer for visualization", 18, False)
    ])
    
    out_path = os.path.join(r"c:\Users\banerjee\Desktop\um1_3d_volume\neurite_detection", "Neurite_Pipeline_Algorithm.pptx")
    prs.save(out_path)
    print(f"Presentation saved successfully at {out_path}")

if __name__ == '__main__':
    main()
