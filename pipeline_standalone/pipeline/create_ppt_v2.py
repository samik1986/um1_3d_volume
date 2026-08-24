import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Neurite Detection Pipeline"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Algorithmic Architecture & Core Routines"
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)

def add_flowchart_slide(prs):
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Pipeline End-to-End Flow"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Stages
    stages = [
        ("1. Input Data", "488nm (Structure)\n555/640nm (Barcodes)"),
        ("2. Feature Extraction", "GPU Soma Masking\nFrangi Skeletonization"),
        ("3. Spatial Filtering", "Isolate True Barcodes\n(+- 5 Voxel Radius)"),
        ("4. Network Assembly", "SWC & JSON Export\nTopological Graphs")
    ]
    
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(2.0)
    height = Inches(2.5)
    arrow_width = Inches(0.3)
    
    for i, (title, desc) in enumerate(stages):
        # Draw Box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(230, 240, 250)
        shape.line.color.rgb = RGBColor(0, 102, 204)
        shape.line.width = Pt(2)
        
        # Add text
        tf = shape.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 51, 102)
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = "\n" + desc
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(50, 50, 50)
        p2.alignment = PP_ALIGN.CENTER
        
        # Draw Arrow
        if i < 3:
            arr_left = left + width + Inches(0.1)
            arr_top = top + height/2 - Inches(0.3)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arr_left, arr_top, arrow_width, Inches(0.6))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(150, 150, 150)
            arrow.line.color.rgb = RGBColor(150, 150, 150)
            
            left = arr_left + arrow_width + Inches(0.1)

def add_detail_slide(prs, title, routines):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for r_title, r_desc in routines:
        p = tf.add_paragraph()
        p.text = r_title + ": "
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)
        
        run = p.add_run()
        run.text = r_desc
        run.font.size = Pt(20)
        run.font.bold = False
        run.font.color.rgb = RGBColor(50, 50, 50)
        
        p.space_before = Pt(14)

def main():
    prs = Presentation()
    
    add_title_slide(prs)
    add_flowchart_slide(prs)
    
    # Slide 2
    add_detail_slide(prs, "Stage 1: Soma Detection (488nm)", [
        ("Parameter Tuning", "np.clip & np.power for Contrast/Gamma limits."),
        ("GPU Smoothing", "cupyx.scipy.ndimage.gaussian_filter applied to raw volume."),
        ("Background Subtraction", "High-radius subtraction to isolate bright cellular regions."),
        ("Thresholding", "cupy.percentile (> 98.5%) to generate binary masks."),
        ("Morphological Clean-up", "cp_ndi.binary_erosion & binary_dilation to smooth borders."),
        ("Size Filtering", "cp.bincount filtering out noise fragments < 3000 voxels.")
    ])

    # Slide 3
    add_detail_slide(prs, "Stage 2: Neurite Skeletonization", [
        ("Vesselness Filtering", "process_tile_frangi_gpu using Multi-scale Analytical Frangi."),
        ("Eigenvalue Solving", "eigh_3x3_analytical via Cardano's exact method on GPU."),
        ("Manifold Thresholding", "cp_ndi.label with dual Hysteresis (High/Low seeds)."),
        ("Soma Subtraction", "Boolean GPU intersection (~gpu_smask) to cut internals."),
        ("Skeletonization", "skimage.morphology.skeletonize using CPU multi-threading.")
    ])

    # Slide 4
    add_detail_slide(prs, "Stage 3: Barcode Structural Filtering", [
        ("Barcode Detection", "detect_barcodes using percentile thresholds (> 99.9%)."),
        ("Centroid Extraction", "scipy.ndimage.center_of_mass on localized spots."),
        ("Spatial Buffer", "cp_ndi.binary_dilation (+-5 voxels) expanding valid skeletons."),
        ("Proximity Check", "filter_structures boolean intersection with expanded mask."),
        ("Mask Finalization", "Retaining only barcodes/neurites within proximity buffer.")
    ])

    # Slide 5
    add_detail_slide(prs, "Stage 4: Network Tracing & Export", [
        ("Subgraph Identification", "scipy.ndimage.label tracing 299,000+ components."),
        ("Path Traversal", "export_graphs recursively walking pixel neighbors."),
        ("Topological JSON", "Constructing CW Complex (nodes/linestrings) for Napari."),
        ("SWC Morphology", "Generating structural tree scaled to physical (Z, Y, X) layout.")
    ])
    
    out_path = os.path.join(r"c:\Users\banerjee\Desktop\um1_3d_volume\neurite_detection", "Neurite_Pipeline_Detailed.pptx")
    prs.save(out_path)
    print(f"Presentation saved successfully at {out_path}")

if __name__ == '__main__':
    main()
