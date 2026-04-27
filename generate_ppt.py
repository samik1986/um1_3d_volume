import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
import zlib
import base64
import urllib.request

def get_flowchart_image(output_path):
    mermaid_text = """%%{init: {'theme': 'base', 'themeVariables': { 'fontSize': '38px', 'fontFamily': 'arial', 'clusterBkg':'#fffbeb', 'clusterBorder':'#d97706'}}}%%
graph LR
    classDef io fill:#fef08a,stroke:#ca8a04,stroke-width:6px,color:#000;
    classDef base fill:#bfdbfe,stroke:#2563eb,stroke-width:6px,color:#000;
    classDef core fill:#86efac,stroke:#059669,stroke-width:6px,color:#000;
    classDef loop fill:#fecaca,stroke:#dc2626,stroke-width:6px,color:#000;

    A(["<b>1. Load FP.tif</b>"]):::io --> B["<b>2. Divide into<br/>4 Quadrants</b>"]:::base
    
    subgraph Quadrant_Loop
        direction LR
        C["<b>3. Compute MIPs &<br/>Save Z-Maps</b>"]:::base --> D["<b>4. Extract Frangi<br/>Skeletons</b>"]:::core
        D --> E["<b>5a. Stereo<br/>Triangulation</b>"]:::core
        D --> F["<b>5b. Vectorized Dense<br/>Back-Project</b>"]:::core
        E --> G["<b>6. Combine Sub-Volume<br/>Mask</b>"]:::base
        F --> G
    end
    
    B --> C
    G --> K{"<b>More<br/>Quadrants?</b>"}:::loop
    K -->|Yes| C
    K -->|No| L["<b>7. Stitch Volume</b>"]:::base
    L --> M(["<b>8. Save Outputs<br/>NRRD / TIFF</b>"]):::io

    style Quadrant_Loop fill:#fffbeb,stroke:#f59e0b,stroke-width:10px
"""
    # Kroki encoding
    compressed = zlib.compress(mermaid_text.encode('utf-8'), 9)
    payload = base64.urlsafe_b64encode(compressed).decode('utf-8')
    url = f"https://kroki.io/mermaid/png/{payload}"
    
    req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
    with urllib.request.urlopen(req) as response, open(output_path, 'wb') as out_file:
        out_file.write(response.read())

def create_presentation():
    prs = pptx.Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "3D Vessel Volume Reconstruction Pipeline"
    subtitle.text = "Results and Step-by-Step Overview\nfp_volume_pipeline"

    # Slide 2: Pipeline Steps
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Pipeline Steps"
    
    tf = body_shape.text_frame
    tf.text = "1. Load FP.tif header (no full RAM load)"
    p = tf.add_paragraph()
    p.text = "2. Divide into 4 XY quadrants (2x2 spatial tiling, all Z slices)"
    p = tf.add_paragraph()
    p.text = "3. Take multi-angle MIP pictures per quadrant (Y-axis camera rotation: 0°, 15°, 30°, 45°) & track Z-Maps"
    p = tf.add_paragraph()
    p.text = "4. Extract vesselness skeletons from each MIP using Frangi/fibermetric logic"
    p = tf.add_paragraph()
    p.text = "5. Back-project skeletons to 3D via:"
    p.level = 0
    p2 = tf.add_paragraph()
    p2.text = "a) Fast Stereo endpoint reconstruction from pairs of 2D views"
    p2.level = 1
    p3 = tf.add_paragraph()
    p3.text = "b) Fully Vectorized Dense Back-Projection using the saved Max Z-Maps"
    p3.level = 1
    p4 = tf.add_paragraph()
    p4.text = "6. Stitch 4 reconstructed quadrants back to the full stitched volume"
    p4.level = 0

    # Slide 3: Flowchart Diagram
    try:
        flowchart_path = os.path.join(os.path.dirname(__file__), "flowchart.png")
        get_flowchart_image(flowchart_path)
        
        slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
        title_shape = slide.shapes.title
        title_shape.text = "Process Flowchart"
        
        # Default PPTX slides are 4:3 (10 inches wide by 7.5 inches high).
        # We'll set the width to 9.0 inches and center it (0.5 inches left margin) so it 
        # fits fully inside the slide without spilling off the edges.
        slide.shapes.add_picture(flowchart_path, Inches(0.5), Inches(2.0), width=Inches(9.0))
    except Exception as e:
        print(f"Flowchart slide skipped due to error: {e}")

    # Slide 4..N: Showing images for each SubVol
    base_dir = r"c:\Users\banerjee\Desktop\um1_3d_volume\fp_pipeline_output"
    
    for vol in range(1, 5):
        # Create a slide for this volume
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
        title_shape = slide.shapes.title
        title_shape.text = f"Quadrant {vol} - Overlays for Angles 0, 15, 30, 45"
        
        angles = ['000', '015', '030', '045']
        
        left_pos = [Inches(0.4), Inches(2.8), Inches(5.2), Inches(7.6)]
        top_pos = Inches(2.5)  # single row
        
        for i, angle in enumerate(angles):
            overlay_file = os.path.join(base_dir, f"SubVol{vol}_Angle{angle}_Overlay.jpg")
            if os.path.exists(overlay_file):
                # Add label above image
                txBox = slide.shapes.add_textbox(left_pos[i], Inches(2.0), Inches(2.2), Inches(0.5))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = f"Angle {int(angle)}°"
                p.alignment = PP_ALIGN.CENTER
                
                # Add picture below label
                slide.shapes.add_picture(overlay_file, left_pos[i], top_pos, width=Inches(2.2))

    # Save presentation
    output_ppt = r"c:\Users\banerjee\Desktop\um1_3d_volume\Pipeline_Results_with_Flowchart.pptx"
    prs.save(output_ppt)
    print(f"Presentation saved to {output_ppt}")

if __name__ == "__main__":
    create_presentation()
